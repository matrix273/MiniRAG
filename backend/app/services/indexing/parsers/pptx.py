"""PowerPoint (.pptx) 文件解析器。"""
from pathlib import Path

from pptx import Presentation

from . import ParseResult


def _table_to_markdown(table) -> str:
    """将 pptx Table 转为 Markdown 表格。"""
    rows_data = []
    for row in table.rows:
        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        rows_data.append(cells)
    if not rows_data:
        return ""
    headers = rows_data[0]
    header_line = "| " + " | ".join(headers) + " |"
    separator = "| " + " | ".join("---" for _ in headers) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows_data[1:])
    return f"{header_line}\n{separator}\n{body}" if body else f"{header_line}\n{separator}"


def parse_pptx(file_path: Path) -> ParseResult:
    """解析 .pptx 文件，每个 Slide 作为一个 section。"""
    prs = Presentation(str(file_path))
    content_parts: list[str] = []
    sections: list[dict] = []
    total_chars = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []

        # 提取所有形状中的文本
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
                        total_chars += len(text)

            if shape.has_table:
                slide_texts.append(_table_to_markdown(shape.table))

        # Speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip()
            if notes_text:
                slide_texts.append(f"**Notes:** {notes_text}")

        if not slide_texts:
            continue

        content_parts.append(f"## Slide {slide_idx}")
        content_parts.extend(slide_texts)

        sections.append({
            "level": 2,
            "title": f"Slide {slide_idx}",
            "text": "\n".join(slide_texts),
        })

    content = "\n\n".join(content_parts)
    metadata = {
        "slide_count": len(prs.slides),
        "total_chars": total_chars,
    }
    return ParseResult(content=content, metadata=metadata, sections=sections)