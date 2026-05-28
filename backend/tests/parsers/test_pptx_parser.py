import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from app.services.indexing.parsers import ParseResult
from app.services.indexing.parsers.pptx import parse_pptx


def _make_pptx(slides_content: list[str]) -> Path:
    """创建测试用 pptx。"""
    prs = Presentation()
    for text in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = text
        slide.placeholders[1].text = f"Content for {text}"
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return Path(tmp.name)


def test_parse_pptx_basic():
    path = _make_pptx(["Slide One", "Slide Two"])
    result = parse_pptx(path)

    assert isinstance(result, ParseResult)
    assert len(result.sections) == 2
    assert result.sections[0]["title"] == "Slide 1"
    assert result.sections[0]["level"] == 2
    assert "## Slide 1" in result.content
    assert "Slide One" in result.content


def test_parse_pptx_metadata():
    path = _make_pptx(["A", "B", "C"])
    result = parse_pptx(path)
    assert result.metadata["slide_count"] == 3


def test_parse_pptx_with_table():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Table Slide"
    # 添加表格
    rows, cols = 2, 2
    left, top, width, height = Inches(1), Inches(1.5), Inches(5), Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "V1"
    table.cell(1, 1).text = "V2"
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)

    result = parse_pptx(Path(tmp.name))
    assert "| H1 |" in result.content
    assert "| V1 |" in result.content
